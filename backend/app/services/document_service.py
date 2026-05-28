from __future__ import annotations

import csv
import json
import mimetypes
import uuid
from dataclasses import dataclass
from datetime import UTC, datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from fastapi import UploadFile

from app.config import Settings, get_settings
from app.core.errors import (
    ConfigurationError,
    DocumentParsingError,
    FeatureDisabledError,
    LLMUnavailableError,
    UnsupportedDocumentError,
)
from app.services.openai_llm_service import OpenAILLMService

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".txt", ".md", ".html", ".htm"}
EXECUTABLE_EXTENSIONS = {
    ".bat",
    ".cmd",
    ".com",
    ".dll",
    ".exe",
    ".js",
    ".msi",
    ".ps1",
    ".scr",
    ".sh",
    ".vbs",
}
ALLOWED_MIME_PREFIXES = {"text/"}
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/markdown",
    "text/html",
    "text/csv",
    "application/csv",
    "application/octet-stream",
}


@dataclass(slots=True)
class ExtractedDocument:
    text: str
    stats: dict[str, Any]
    metadata: dict[str, Any]


class DocumentService:
    def __init__(
        self,
        settings: Settings | None = None,
        llm: OpenAILLMService | None = None,
    ) -> None:
        self._settings = settings or get_settings()
        self._llm = llm or OpenAILLMService()

    async def upload(self, *, file: UploadFile, user_id: str | None = None) -> dict[str, Any]:
        self._ensure_enabled()
        filename = Path(file.filename or "document").name
        extension = Path(filename).suffix.lower()
        content_type = file.content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
        self._validate_file(filename=filename, extension=extension, content_type=content_type)

        document_id = str(uuid.uuid4())
        created_at = datetime.now(UTC).isoformat()
        upload_dir = self._upload_dir()
        self._cleanup_expired_uploads(upload_dir)
        target = upload_dir / f"{document_id}{extension}"

        size = 0
        with target.open("wb") as handle:
            while chunk := await file.read(1024 * 1024):
                size += len(chunk)
                if size > self._settings.document_max_upload_bytes:
                    target.unlink(missing_ok=True)
                    raise ConfigurationError(
                        f"File is too large. Maximum size is {self._settings.document_max_upload_bytes} bytes."
                    )
                handle.write(chunk)

        metadata = {
            "documentId": document_id,
            "filename": filename,
            "contentType": content_type,
            "sizeBytes": size,
            "extension": extension,
            "createdAt": created_at,
            "userId": user_id,
            "retentionHours": self._settings.document_retention_hours,
        }
        self._metadata_path(document_id).write_text(json.dumps(metadata), encoding="utf-8")
        return {
            "documentId": document_id,
            "filename": filename,
            "contentType": content_type,
            "sizeBytes": size,
            "status": "uploaded",
            "createdAt": created_at,
        }

    async def analyze(
        self, *, document_id: str, analysis_type: str, instructions: str
    ) -> dict[str, Any]:
        self._ensure_enabled()
        extracted = self.extract(document_id)
        try:
            result = await self._llm.complete_json(
                system_prompt=_analysis_system_prompt(),
                user_prompt=_analysis_user_prompt(
                    extracted=extracted,
                    analysis_type=analysis_type,
                    instructions=instructions,
                ),
            )
        except Exception as exc:
            raise LLMUnavailableError("Could not analyze document because the LLM is unavailable.") from exc

        return {
            "documentId": document_id,
            "filename": extracted.metadata["filename"],
            "analysisType": analysis_type,
            "summary": str(result.get("summary") or ""),
            "keyPoints": _string_list(result.get("keyPoints")),
            "risks": _string_list(result.get("risks")),
            "actionItems": _string_list(result.get("actionItems")),
            "confidence": _confidence(result.get("confidence")),
            "sourceStats": extracted.stats,
            "llmStatus": "ok",
        }

    async def report(
        self,
        *,
        document_id: str,
        report_title: str,
        instructions: str,
        report_format: str,
    ) -> dict[str, Any]:
        self._ensure_enabled()
        extracted = self.extract(document_id)
        try:
            result = await self._llm.complete_json(
                system_prompt=_report_system_prompt(),
                user_prompt=_report_user_prompt(
                    extracted=extracted,
                    title=report_title,
                    instructions=instructions,
                    report_format=report_format,
                ),
            )
        except Exception as exc:
            raise LLMUnavailableError("Could not generate report because the LLM is unavailable.") from exc

        sections = result.get("sections")
        if not isinstance(sections, list):
            sections = []
        normalized_sections = [
            {
                "heading": str(section.get("heading") or "Section"),
                "content": str(section.get("content") or ""),
            }
            for section in sections
            if isinstance(section, dict)
        ]
        created_at = datetime.now(UTC).isoformat()
        return {
            "documentId": document_id,
            "reportId": str(uuid.uuid4()),
            "filename": extracted.metadata["filename"],
            "title": str(result.get("title") or report_title),
            "report": str(result.get("report") or ""),
            "sections": normalized_sections,
            "sourceStats": extracted.stats,
            "llmStatus": "ok",
            "createdAt": created_at,
        }

    def extract(self, document_id: str) -> ExtractedDocument:
        metadata = self._load_metadata(document_id)
        path = self._document_path(document_id, metadata["extension"])
        if not path.exists() or not path.is_file():
            raise DocumentParsingError("Uploaded document was not found.")

        extension = metadata["extension"]
        if extension == ".pdf":
            text, stats = _extract_pdf(path)
        elif extension == ".docx":
            text, stats = _extract_docx(path)
        elif extension == ".pptx":
            text, stats = _extract_pptx(path)
        elif extension == ".xlsx":
            text, stats = _extract_xlsx(path)
        elif extension == ".csv":
            text, stats = _extract_csv(path)
        elif extension in {".txt", ".md", ".html", ".htm"}:
            text, stats = _extract_txt(path)
        else:
            raise UnsupportedDocumentError(f"Unsupported file extension: {extension}")

        text = _strip_html(text).strip()
        if not text:
            raise DocumentParsingError("No readable text could be extracted from this document.")
        return ExtractedDocument(
            text=text[:60_000],
            stats={
                **stats,
                "charactersExtracted": len(text),
                "truncatedForLlm": len(text) > 60_000,
            },
            metadata=metadata,
        )

    def _validate_file(self, *, filename: str, extension: str, content_type: str) -> None:
        if extension in EXECUTABLE_EXTENSIONS:
            raise UnsupportedDocumentError(f"Executable files are not supported: {filename}.")
        if extension not in ALLOWED_EXTENSIONS:
            raise UnsupportedDocumentError(f"Unsupported file type for {filename}.")
        if not (
            content_type in ALLOWED_MIME_TYPES
            or any(content_type.startswith(prefix) for prefix in ALLOWED_MIME_PREFIXES)
        ):
            raise UnsupportedDocumentError(f"Unsupported MIME type: {content_type}.")

    def _upload_dir(self) -> Path:
        path = Path(self._settings.document_upload_dir)
        if not path.is_absolute():
            path = Path.cwd() / path
        path.mkdir(parents=True, exist_ok=True)
        return path.resolve()

    def _metadata_path(self, document_id: str) -> Path:
        return self._upload_dir() / f"{document_id}.json"

    def _document_path(self, document_id: str, extension: str) -> Path:
        return self._upload_dir() / f"{document_id}{extension}"

    def _load_metadata(self, document_id: str) -> dict[str, Any]:
        path = self._metadata_path(document_id)
        if not path.exists() or not path.is_file():
            raise DocumentParsingError("Uploaded document metadata was not found.")
        payload = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            raise DocumentParsingError("Uploaded document metadata is invalid.")
        return payload

    def _ensure_enabled(self) -> None:
        if not self._settings.enable_real_document_analysis:
            raise FeatureDisabledError(
                "Real document analysis is disabled. Set ENABLE_REAL_DOCUMENT_ANALYSIS=true to enable it."
            )

    def _cleanup_expired_uploads(self, upload_dir: Path) -> None:
        retention_hours = self._settings.document_retention_hours
        if retention_hours <= 0:
            return
        now = datetime.now(UTC)
        for metadata_path in upload_dir.glob("*.json"):
            try:
                metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
                if not isinstance(metadata, dict):
                    continue
                created_at = _parse_datetime(str(metadata.get("createdAt") or ""))
                age_hours = (now - created_at).total_seconds() / 3600
                if age_hours <= retention_hours:
                    continue
                extension = str(metadata.get("extension") or "")
                document_id = metadata_path.stem
                candidates = [metadata_path]
                if extension:
                    candidates.append(self._document_path(document_id, extension))
                for candidate in candidates:
                    resolved = candidate.resolve()
                    if resolved.is_relative_to(upload_dir.resolve()):
                        resolved.unlink(missing_ok=True)
            except Exception:
                continue


def _extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentParsingError("PDF parsing dependency pypdf is not installed.") from exc
    try:
        reader = PdfReader(str(path))
        texts = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise DocumentParsingError("Could not parse PDF document.") from exc
    return "\n\n".join(texts), {"pages": len(reader.pages), "sheets": 0, "parser": "pypdf"}


def _extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentParsingError("DOCX parsing dependency python-docx is not installed.") from exc
    try:
        document = Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        table_cells = [
            cell.text
            for table in document.tables
            for row in table.rows
            for cell in row.cells
            if cell.text
        ]
    except Exception as exc:
        raise DocumentParsingError("Could not parse DOCX document.") from exc
    return "\n".join([*paragraphs, *table_cells]), {"pages": None, "sheets": 0, "parser": "python-docx"}


def _extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentParsingError("PPTX parsing dependency python-pptx is not installed.") from exc
    try:
        presentation = Presentation(str(path))
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text.strip():
                        slide_lines.append(text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text and cell.text.strip()
                        ]
                        if values:
                            slide_lines.append(" | ".join(values))
            if slide_lines:
                lines.append(f"Slide {index}\n" + "\n".join(slide_lines))
    except Exception as exc:
        raise DocumentParsingError("Could not parse PPTX document.") from exc
    return "\n\n".join(lines), {"pages": len(presentation.slides), "sheets": 0, "parser": "python-pptx"}


def _extract_xlsx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise DocumentParsingError("XLSX parsing dependency openpyxl is not installed.") from exc
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    lines.append(" | ".join(values))
    except Exception as exc:
        raise DocumentParsingError("Could not parse XLSX document.") from exc
    return "\n".join(lines), {"pages": None, "sheets": len(workbook.worksheets), "parser": "openpyxl"}


def _extract_csv(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
    except UnicodeDecodeError:
        with path.open("r", encoding="latin-1", newline="") as handle:
            rows = list(csv.reader(handle))
    except Exception as exc:
        raise DocumentParsingError("Could not parse CSV document.") from exc
    return "\n".join(" | ".join(row) for row in rows), {"pages": None, "sheets": 1, "rows": len(rows), "parser": "csv"}


def _extract_txt(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        return path.read_text(encoding="utf-8"), {"pages": None, "sheets": 0, "parser": "text"}
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1"), {"pages": None, "sheets": 0, "parser": "text"}
    except Exception as exc:
        raise DocumentParsingError("Could not parse TXT document.") from exc


class _HTMLTextStripper(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        self.parts.append(data)


def _strip_html(value: str) -> str:
    parser = _HTMLTextStripper()
    parser.feed(value)
    stripped = " ".join(part.strip() for part in parser.parts if part.strip())
    return stripped or value


def _parse_datetime(value: str) -> datetime:
    parsed = datetime.fromisoformat(value)
    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=UTC)
    return parsed.astimezone(UTC)


def _analysis_system_prompt() -> str:
    return """You analyze uploaded business documents for NexusHub.
Return strict JSON only: summary, keyPoints, risks, actionItems, confidence.
Use only extracted document content. Do not invent facts. If content is insufficient, say so clearly.
Produce executive-friendly summaries with risks, decisions, action items, and deadlines when available.
Avoid generic filler and never say "as an AI"."""


def _analysis_user_prompt(*, extracted: ExtractedDocument, analysis_type: str, instructions: str) -> str:
    return f"""Filename: {extracted.metadata['filename']}
Analysis type: {analysis_type}
Instructions: {instructions or "Create an executive brief."}

Extracted content:
{extracted.text}"""


def _report_system_prompt() -> str:
    return """You generate concise executive reports from uploaded documents.
Return strict JSON only: title, report, sections.
sections must be an array of objects with heading and content.
Use only extracted document content. Do not invent facts, fake dates, or fake commitments.
If content is insufficient, state the limitation in the report.
Avoid generic filler and never say "as an AI"."""


def _report_user_prompt(
    *,
    extracted: ExtractedDocument,
    title: str,
    instructions: str,
    report_format: str,
) -> str:
    return f"""Requested title: {title}
Format: {report_format}
Instructions: {instructions or "Generate an executive summary report."}

Extracted content:
{extracted.text}"""


def _string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value if str(item).strip()]
    if isinstance(value, str) and value.strip():
        return [value]
    return []


def _confidence(value: Any) -> float:
    try:
        parsed = float(value)
    except (TypeError, ValueError):
        return 0.5
    return max(0.0, min(parsed, 1.0))
